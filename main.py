from fastapi import FastAPI, HTTPException, Depends, UploadFile, File, Form
from fastapi.staticfiles import StaticFiles
from fastapi.responses import HTMLResponse, StreamingResponse, Response
from fastapi.middleware.cors import CORSMiddleware
from sqlalchemy.orm import Session
from pydantic import BaseModel
from typing import Optional, List
import csv
import io
import json
import zipfile
from urllib.parse import quote
from datetime import datetime, timezone

from database import get_db, engine
import models
from models import JobPosting, Setting
from ai_service import generate_job_posting
from bulk_convert import convert_jobbudy_to_indeed
from atally_convert import convert_jobbudy_to_atally

models.Base.metadata.create_all(bind=engine)

# 既存DBへの列追加マイグレーション
from sqlalchemy import text as _text
with engine.connect() as _conn:
    for _col in ["hire_count INTEGER", "haken_company_name VARCHAR(200)", "haken_address VARCHAR(500)", "haken_notes TEXT", "job_category VARCHAR(100)"]:
        try:
            _conn.execute(_text(f"ALTER TABLE job_postings ADD COLUMN {_col}"))
            _conn.commit()
        except Exception:
            pass

app = FastAPI(title="Indeed求人作る君", version="1.0.0")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

# ローカル開発時はFastAPIが静的ファイルを配信
# PythonAnywhereではダッシュボードのStatic Filesで /static/ を設定するため不要だが
# フォールバックとして残しておく（SERVE_STATIC=false で無効化可能）
import os as _os
if _os.getenv("SERVE_STATIC", "true").lower() != "false":
    app.mount("/static", StaticFiles(directory="static"), name="static")


# ---------- Pydantic schemas ----------


class GenerateRequest(BaseModel):
    company_url: str
    request_text: str
    recruitment_url: str = ""
    application_url: str = ""
    contact_name: str = ""
    contact_phone: str = ""
    contact_email: str = ""
    target_persona: str = ""
    employment_type: str = ""


class JobCreate(BaseModel):
    company_name: str
    company_url: Optional[str] = None
    job_title: str
    prefecture: Optional[str] = None
    city: Optional[str] = None
    employment_type: Optional[str] = None
    salary_min: Optional[int] = None
    salary_max: Optional[int] = None
    salary_type: Optional[str] = None
    description: Optional[str] = None
    requirements: Optional[str] = None
    preferred_skills: Optional[str] = None
    working_hours: Optional[str] = None
    holidays: Optional[str] = None
    benefits: Optional[str] = None
    selection_process: Optional[str] = None
    hire_count: Optional[int] = None
    haken_company_name: Optional[str] = None
    haken_address: Optional[str] = None
    haken_notes: Optional[str] = None
    job_category: Optional[str] = None
    appeal_points: Optional[str] = None
    application_url: Optional[str] = None
    contact_name: Optional[str] = None
    contact_phone: Optional[str] = None
    contact_email: Optional[str] = None
    original_request: Optional[str] = None


class JobUpdate(JobCreate):
    company_name: Optional[str] = None
    job_title: Optional[str] = None


class ExportRequest(BaseModel):
    job_ids: List[int]


# ---------- helpers ----------


def job_to_dict(job: JobPosting) -> dict:
    return {
        "id": job.id,
        "company_name": job.company_name,
        "company_url": job.company_url,
        "job_title": job.job_title,
        "prefecture": job.prefecture,
        "city": job.city,
        "employment_type": job.employment_type,
        "salary_min": job.salary_min,
        "salary_max": job.salary_max,
        "salary_type": job.salary_type,
        "description": job.description,
        "requirements": job.requirements,
        "preferred_skills": job.preferred_skills,
        "working_hours": job.working_hours,
        "holidays": job.holidays,
        "benefits": job.benefits,
        "selection_process": job.selection_process,
        "hire_count": job.hire_count,
        "haken_company_name": job.haken_company_name,
        "haken_address": job.haken_address,
        "haken_notes": job.haken_notes,
        "job_category": job.job_category,
        "appeal_points": job.appeal_points,
        "application_url": job.application_url,
        "contact_name": job.contact_name,
        "contact_phone": job.contact_phone,
        "contact_email": job.contact_email,
        "original_request": job.original_request,
        "created_at": job.created_at.isoformat() if job.created_at else None,
        "updated_at": job.updated_at.isoformat() if job.updated_at else None,
    }


# ---------- routes ----------


def _extract_text_from_file(filename: str, content: bytes) -> str:
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""

    if ext == "pdf":
        from pypdf import PdfReader
        reader = PdfReader(io.BytesIO(content))
        return "\n".join(
            page.extract_text() or "" for page in reader.pages
        ).strip()

    if ext == "docx":
        from docx import Document
        doc = Document(io.BytesIO(content))
        return "\n".join(p.text for p in doc.paragraphs if p.text).strip()

    if ext == "pptx":
        from pptx import Presentation
        prs = Presentation(io.BytesIO(content))
        lines = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            lines.append(text)
        return "\n".join(lines)

    if ext in ("txt", "csv", "md"):
        for enc in ("utf-8", "shift_jis", "cp932", "euc_jp"):
            try:
                return content.decode(enc).strip()
            except UnicodeDecodeError:
                continue
        return content.decode("utf-8", errors="replace").strip()

    raise ValueError(f"非対応のファイル形式です: .{ext}")


@app.post("/api/extract-text")
async def extract_text(file: UploadFile = File(...)):
    if file.size and file.size > 20 * 1024 * 1024:
        raise HTTPException(status_code=400, detail="ファイルサイズは20MB以下にしてください")
    content = await file.read()
    try:
        text = _extract_text_from_file(file.filename or "", content)
    except ValueError as e:
        raise HTTPException(status_code=400, detail=str(e))
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"テキスト抽出に失敗しました: {e}")
    if not text:
        raise HTTPException(status_code=400, detail="ファイルからテキストを抽出できませんでした")
    return {"text": text, "filename": file.filename}


@app.post("/api/bulk-convert")
async def bulk_convert(file: UploadFile = File(...), target: str = Form("indeed")):
    """Jobbudyの求人一覧Excelを取り込み、Indeed形式またはAtally形式に一括変換。
    件数ごとに分割し、複数ファイルになる場合はZIPで返す。
    target: "indeed"（既定）または "atally"。"""
    name = file.filename or ""
    if not name.lower().endswith((".xlsx", ".xls")):
        raise HTTPException(status_code=400, detail="Excelファイル（.xlsx / .xls）を添付してください")
    if file.size and file.size > 150 * 1024 * 1024:
        raise HTTPException(status_code=400, detail="ファイルサイズは150MB以下にしてください")

    target = (target or "indeed").lower()
    if target not in ("indeed", "atally"):
        raise HTTPException(status_code=400, detail="target は indeed または atally を指定してください")

    content = await file.read()
    try:
        if target == "atally":
            files, stats = convert_jobbudy_to_atally(content, source_filename=name)
        else:
            files, stats = convert_jobbudy_to_indeed(content, source_filename=name)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"変換に失敗しました: {e}")

    if not files or stats["converted"] == 0:
        raise HTTPException(status_code=400, detail="変換できる求人がありませんでした")

    prefix = "atally_jobs" if target == "atally" else "indeed_jobs"

    stats_header = {
        "X-Convert-Total": str(stats["total"]),
        "X-Convert-Converted": str(stats["converted"]),
        "X-Convert-Skipped": str(stats["skipped"]),
        "X-Convert-Parts": str(stats["parts"]),
        "Access-Control-Expose-Headers": "Content-Disposition,X-Convert-Total,X-Convert-Converted,X-Convert-Skipped,X-Convert-Parts",
    }

    # 1ファイルのみならCSVをそのまま、複数ならZIPにまとめて返す
    if len(files) == 1:
        fname, data = files[0]
        # 日本語ファイル名はRFC 5987でエンコード（HTTPヘッダはlatin-1のため）
        disposition = f"attachment; filename={prefix}.csv; filename*=UTF-8''{quote(fname)}"
        return Response(
            content=data,
            media_type="text/csv; charset=utf-8",
            headers={"Content-Disposition": disposition, **stats_header},
        )

    zip_buf = io.BytesIO()
    with zipfile.ZipFile(zip_buf, "w", zipfile.ZIP_DEFLATED) as zf:
        for fname, data in files:
            zf.writestr(fname, data)
    zip_name = f"{prefix}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.zip"
    return Response(
        content=zip_buf.getvalue(),
        media_type="application/zip",
        headers={"Content-Disposition": f"attachment; filename={zip_name}", **stats_header},
    )


@app.get("/api/settings")
def get_settings(db: Session = Depends(get_db)):
    rows = db.query(Setting).all()
    return {r.key: r.value for r in rows}


@app.put("/api/settings")
def update_settings(data: dict, db: Session = Depends(get_db)):
    for key, value in data.items():
        row = db.query(Setting).filter(Setting.key == key).first()
        if row:
            row.value = value
        else:
            db.add(Setting(key=key, value=value))
    db.commit()
    return {"ok": True}


@app.get("/", response_class=HTMLResponse)
async def root():
    with open("static/index.html", "r", encoding="utf-8") as f:
        return HTMLResponse(f.read())


@app.post("/api/generate")
def generate(req: GenerateRequest, db: Session = Depends(get_db)):
    if not req.company_url.startswith("http"):
        raise HTTPException(status_code=400, detail="企業URLはhttp/httpsで始まる必要があります")
    if not req.request_text.strip():
        raise HTTPException(status_code=400, detail="依頼文を入力してください")
    try:
        haken_name = ""
        if req.employment_type == "派遣社員":
            haken_row = db.query(Setting).filter(Setting.key == "haken_company_name").first()
            haken_name = haken_row.value if haken_row else ""
        result = generate_job_posting(
            company_url=req.company_url,
            request_text=req.request_text,
            recruitment_url=req.recruitment_url,
            application_url=req.application_url,
            contact_name=req.contact_name,
            contact_phone=req.contact_phone,
            contact_email=req.contact_email,
            target_persona=req.target_persona,
            employment_type=req.employment_type,
            haken_company_name=haken_name,
        )
        return result
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.get("/api/jobs")
def list_jobs(db: Session = Depends(get_db)):
    jobs = db.query(JobPosting).order_by(JobPosting.created_at.desc()).all()
    return [job_to_dict(j) for j in jobs]


@app.post("/api/jobs")
def create_job(job_data: JobCreate, db: Session = Depends(get_db)):
    job = JobPosting(**job_data.model_dump())
    db.add(job)
    db.commit()
    db.refresh(job)
    return job_to_dict(job)


@app.get("/api/jobs/{job_id}")
def get_job(job_id: int, db: Session = Depends(get_db)):
    job = db.query(JobPosting).filter(JobPosting.id == job_id).first()
    if not job:
        raise HTTPException(status_code=404, detail="求人が見つかりません")
    return job_to_dict(job)


@app.put("/api/jobs/{job_id}")
def update_job(job_id: int, job_data: JobUpdate, db: Session = Depends(get_db)):
    job = db.query(JobPosting).filter(JobPosting.id == job_id).first()
    if not job:
        raise HTTPException(status_code=404, detail="求人が見つかりません")
    for key, value in job_data.model_dump(exclude_none=True).items():
        setattr(job, key, value)
    job.updated_at = datetime.now(timezone.utc)
    db.commit()
    db.refresh(job)
    return job_to_dict(job)


@app.delete("/api/jobs/{job_id}")
def delete_job(job_id: int, db: Session = Depends(get_db)):
    job = db.query(JobPosting).filter(JobPosting.id == job_id).first()
    if not job:
        raise HTTPException(status_code=404, detail="求人が見つかりません")
    db.delete(job)
    db.commit()
    return {"message": "削除しました"}


def _build_other(job) -> str:
    parts = []
    if job.selection_process:
        parts.append(job.selection_process)
    if job.employment_type == "派遣社員":
        haken_lines = []
        if job.haken_company_name:
            haken_lines.append(f"派遣元事業者名：{job.haken_company_name}")
        if job.haken_address:
            haken_lines.append(f"派遣元住所：{job.haken_address}")
        if job.haken_notes:
            haken_lines.append(job.haken_notes)
        if haken_lines:
            parts.append("\n".join(haken_lines))
    return "\n\n".join(parts)


@app.post("/api/export")
def export_jobs(req: ExportRequest, db: Session = Depends(get_db)):
    jobs = db.query(JobPosting).filter(JobPosting.id.in_(req.job_ids)).all()
    if not jobs:
        raise HTTPException(status_code=404, detail="エクスポート対象の求人が見つかりません")

    output = io.StringIO()
    writer = csv.writer(output)

    # Indeed Japan 公式フォーマット（69列・テンプレート完全準拠）
    headers = [
        "ステータス", "会社名", "職種名", "職業カテゴリー", "求人キャッチコピー",
        "勤務地（郵便番号）", "勤務地（都道府県・市区町村・町域）", "勤務地（丁目・番地・号）", "勤務地（建物名・階数）",
        "雇用形態", "有料職業紹介に該当",
        "給与形態", "給与（最低額）", "給与（最高額）", "給与（表示形式）",
        "固定残業代の有無", "固定残業代（最低額）", "固定残業代（最高額）", "固定残業代（支払い単位）",
        "固定残業代（時間）", "固定残業代（分）", "固定残業代（超過分の追加支払への同意）",
        "勤務形態", "平均所定労働時間", "平均所定労働時間（分）",
        "社会保険", "社会保険（適用されない理由）",
        "試用期間の有無", "試用期間（期間）", "試用期間（期間の単位）", "試用期間（試用期間中の労働条件）",
        "試用期間中の給与形態", "試用期間中の給与（最低額）", "試用期間中の給与（最高額）", "試用期間中の給与（表示形式）",
        "試用期間中の固定残業代の有無", "試用期間中の固定残業代（最低額）", "試用期間中の固定残業代（最高額）",
        "試用期間中の固定残業代（支払い単位）", "試用期間中の固定残業代（時間）", "試用期間中の固定残業代（分）",
        "試用期間中の固定残業代（超過分の追加支払への同意）",
        "試用期間中の平均所定労働時間", "試用期間中の平均所定労働時間（分）", "試用期間中のその他の条件",
        "募集要項（仕事内容）", "募集要項（アピールポイント）", "募集要項（求める人材）",
        "募集要項（勤務時間・曜日）", "募集要項（休暇・休日）", "募集要項（勤務地の補足）",
        "募集要項（アクセス）", "募集要項（給与の補足）", "募集要項（待遇・福利厚生）", "募集要項（その他）",
        "掲載画像", "タグ", "タグ", "タグ",
        "採用予定人数", "履歴書の有無",
        "応募者に関する情報", "応募者に関する情報", "応募者に関する情報",
        "応募用メールアドレス", "求人問い合わせ先電話番号（半角）",
        "審査用の質問", "自動アプローチ利用設定", "ユーザー指定ID",
    ]
    writer.writerow(headers)

    for job in jobs:
        # アピールポイントを取得
        appeal_pts = []
        if job.appeal_points:
            try:
                appeal_pts = json.loads(job.appeal_points)
            except Exception:
                appeal_pts = [job.appeal_points]
        catchcopy = appeal_pts[0] if appeal_pts else ""
        appeal_text = "\n".join(appeal_pts)

        # 求める人材（必須＋歓迎）
        requirements_combined = "\n".join(filter(None, [
            job.requirements or "",
            f"【歓迎】\n{job.preferred_skills}" if job.preferred_skills else "",
        ]))

        # 勤務地
        location = "".join(filter(None, [job.prefecture or "", job.city or ""]))

        # 給与の最低＜最高を保証
        sal_min = job.salary_min
        sal_max = job.salary_max
        if sal_min and sal_max:
            if sal_min > sal_max:
                sal_min, sal_max = sal_max, sal_min
            if sal_min == sal_max:  # 同額はIndeedが拒否するため最高額を削除
                sal_max = None

        # 旧形式の職業カテゴリーを有効値にマッピング
        _CAT_MAP = {
            "クリエイティブ": "その他クリエイティブ/デザイン職",
            "事務・管理": "一般事務",
            "IT・機械・電気・電子": "その他IT",
            "建設・土木・設備・設計": "その他建築/土木/プラント専門職",
            "医療・介護・福祉": "医療/福祉専門職その他",
            "飲食・フード": "その他飲食",
            "製造・工場・物流・倉庫": "その他機械/電気/電子製品専門職",
            "保育・教育": "その他教育/保育専門職",
            "農林水産": "その他公務員/団体職員/農林水産",
            "販売・サービス": "その他販売/フロアスタッフ",
            "企画・マーケティング・経営・管理職": "その他企画/マーケティング",
        }
        job_category_csv = job.job_category or ""
        job_category_csv = _CAT_MAP.get(job_category_csv, job_category_csv)
        if not job_category_csv:
            job_category_csv = "一般事務"  # 未設定レコードのフォールバック

        row = [
            "募集中",                              # ステータス
            job.company_name or "",               # 会社名
            job.job_title or "",                  # 職種名
            job_category_csv,                     # 職業カテゴリー
            catchcopy,                            # 求人キャッチコピー
            "",                                   # 勤務地（郵便番号）
            location,                             # 勤務地（都道府県・市区町村・町域）
            "",                                   # 勤務地（丁目・番地・号）
            "",                                   # 勤務地（建物名・階数）
            job.employment_type or "",            # 雇用形態
            "いいえ",                              # 有料職業紹介に該当
            job.salary_type or "",                # 給与形態
            sal_min or "",                        # 給与（最低額）
            sal_max or "",                        # 給与（最高額）
            "範囲で表示" if (sal_min and sal_max) else ("最低額を表示" if sal_min else "固定額を表示"),  # 給与（表示形式）
            "", "", "", "", "", "", "",           # 固定残業代関連（7列）
            "固定時間制",                          # 勤務形態
            "", "",                               # 平均所定労働時間
            "健康保険,厚生年金,雇用保険,労災保険", "",  # 社会保険 / 適用されない理由
            "なし", "", "", "",                   # 試用期間（有無・期間・単位・条件）
            "", "", "", "",                       # 試用期間中の給与
            "", "", "", "", "", "", "",           # 試用期間中の固定残業代（7列）
            "", "",                               # 試用期間中の平均所定労働時間
            "",                                   # 試用期間中のその他の条件
            job.description or "",           # 募集要項（仕事内容）
            appeal_text,                     # 募集要項（アピールポイント）
            requirements_combined,           # 募集要項（求める人材）
            job.working_hours or "",         # 募集要項（勤務時間・曜日）
            job.holidays or "",              # 募集要項（休暇・休日）
            "",                              # 募集要項（勤務地の補足）
            "",                              # 募集要項（アクセス）
            "",                              # 募集要項（給与の補足）
            job.benefits or "",              # 募集要項（待遇・福利厚生）
            _build_other(job),               # 募集要項（その他）
            "",                              # 掲載画像
            "", "", "",                      # タグ（3列）
            job.hire_count or "",            # 採用予定人数
            "必須",                           # 履歴書の有無
            "", "", "",                      # 応募者に関する情報（3列）
            job.contact_email or "nishihara@toyopla.jp",  # 応募用メールアドレス
            job.contact_phone or "",         # 求人問い合わせ先電話番号
            "", "",                          # 審査用の質問・自動アプローチ利用設定
            "",                              # ユーザー指定ID
        ]
        writer.writerow(row)

    # UTF-8 with BOM（Excel・Indeed Japan どちらでも正しく読める）
    csv_bytes = b'\xef\xbb\xbf' + output.getvalue().encode("utf-8")
    filename = f"indeed_jobs_{datetime.now().strftime('%Y%m%d_%H%M%S')}.csv"

    return Response(
        content=csv_bytes,
        media_type="text/csv; charset=utf-8",
        headers={"Content-Disposition": f"attachment; filename={filename}"},
    )
